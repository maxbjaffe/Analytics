import streamlit as st
import pandas as pd
import matplotlib.pyplot as plt
import matplotlib.ticker as mtick
import io
from reportlab.pdfgen import canvas
from pptx import Presentation
from pptx.util import Inches

# === Page Setup ===
st.set_page_config(page_title="TVQI Dashboard", layout="wide")

# === Header Background Styling ===
st.markdown(
    """
    <style>
    .header-container {
        position: relative;
        height: 200px;
        background-image: url('dashboard_header.png');
        background-size: cover;
        background-position: center;
        opacity: 0.3;
        margin-bottom: -120px;
    }
    .title-text {
        position: relative;
        text-align: center;
        font-size: 40px;
        font-weight: bold;
        padding-top: 40px;
        z-index: 2;
    }
    </style>
    <div class="header-container"></div>
    <div class="title-text">📊 TV Quality Index Dashboard</div>
    """,
    unsafe_allow_html=True
)

# === Utility Functions ===
def safe_format_number(val, as_int=True):
    try:
        return f"{int(val):,}" if as_int else f"{float(val):,.2f}"
    except (ValueError, TypeError):
        return "N/A"

def apply_mapping(df, column, path):
    try:
        mapping_df = pd.read_csv(path)
        name_map = dict(zip(mapping_df["raw_name"], mapping_df["clean_name"]))
        df[column] = df[column].replace(name_map)
    except Exception as e:
        st.warning(f"⚠️ Could not apply mapping for '{column}': {e}")
    return df

# === File Upload ===
uploaded_file = st.file_uploader("Upload your TVQI Report CSV", type=["csv"])

if uploaded_file:
    df = pd.read_csv(uploaded_file)
    df.columns = [col.strip().lower().replace(" ", "_") for col in df.columns]
    df.dropna(how="all", inplace=True)

    # Apply mappings
    df = apply_mapping(df, "supply_vendor", "vendor_mapping.csv")
    df = apply_mapping(df, "campaign", "campaign_mapping.csv")
    df = apply_mapping(df, "inventory_contract", "inventory_mapping.csv")

    df["advertiser_cost_(adv_currency)"] = pd.to_numeric(df["advertiser_cost_(adv_currency)"], errors="coerce")
    df["tvqi_score"] = df["tv_quality_index_raw"] / df["tv_quality_index_measured_impressions"]

    st.subheader("📁 Uploaded Data")
    st.dataframe(df)

    # === Summary Metrics ===
    st.subheader("📈 Summary Metrics")
    impressions = df["tv_quality_index_measured_impressions"].sum()
    completed_views = df["player_completed_views"].sum()
    player_starts = df["player_starts"].sum()
    viewable_impressions = df["sampled_viewed_impressions"].sum()
    cost = df["advertiser_cost_(adv_currency)"].sum()
    tvqi_raw = df["tv_quality_index_raw"].sum()
    tvqi_score = (tvqi_raw / impressions) if impressions > 0 else None
    cpm = (cost / impressions * 1000) if impressions > 0 else None
    completion_rate = (completed_views / player_starts) if player_starts > 0 else None
    ecpcv = (cost / completed_views) if completed_views > 0 else None
    viewable_cpm = (cost / viewable_impressions * 1000) if viewable_impressions > 0 else None
    in_view_rate = (viewable_impressions / impressions) if impressions > 0 else None

    col1, col2, col3 = st.columns(3)
    col1.metric("TVQI Measured Impressions", safe_format_number(impressions))
    col2.metric("Completed Views", safe_format_number(completed_views))
    col3.metric("Total Cost", f"${safe_format_number(cost, as_int=False)}")

    col4, col5, col6 = st.columns(3)
    col4.metric("CPM", f"${safe_format_number(cpm, as_int=False)}")
    col5.metric("Completion Rate", f"{completion_rate:.2%}" if completion_rate else "N/A")
    col6.metric("eCPCV", f"${safe_format_number(ecpcv, as_int=False)}")

    col7, col8, col9 = st.columns(3)
    col7.metric("Viewable CPM", f"${safe_format_number(viewable_cpm, as_int=False)}")
    col8.metric("In-View Rate", f"{in_view_rate:.2%}" if in_view_rate else "N/A")
    col9.metric("TVQI Score", f"{tvqi_score:.4f}" if tvqi_score else "N/A")

    # === Charts ===
    st.subheader("📊 Top 10: Supply Vendor vs Inventory Contract")
    col_chart1, col_chart2 = st.columns(2)

    # --- Supply Vendor Chart ---
    with col_chart1:
        vendor_group = df.groupby("supply_vendor").agg({
            "advertiser_cost_(adv_currency)": "sum",
            "tv_quality_index_raw": "sum",
            "tv_quality_index_measured_impressions": "sum"
        }).reset_index()
        vendor_group["tvqi_score"] = vendor_group["tv_quality_index_raw"] / vendor_group["tv_quality_index_measured_impressions"]
        vendor_group = vendor_group.sort_values(by="advertiser_cost_(adv_currency)", ascending=False).head(10)

        fig1, ax1 = plt.subplots(figsize=(7, 3.5))
        ax2 = ax1.twinx()
        ax1.bar(vendor_group["supply_vendor"], vendor_group["advertiser_cost_(adv_currency)"], color="lightgreen")
        ax2.plot(vendor_group["supply_vendor"], vendor_group["tvqi_score"], color="blue", marker="o")
        ax1.set_ylabel("Advertiser Cost ($)", color="green")
        ax2.set_ylabel("TVQI Score", color="blue")
        ax1.set_title("By Supply Vendor", fontsize=12)
        ax1.tick_params(axis="x", rotation=45)
        ax1.yaxis.set_major_formatter(mtick.StrMethodFormatter("${x:,.0f}"))
        fig1.tight_layout()
        st.pyplot(fig1)

    # --- Inventory Contract Chart ---
    with col_chart2:
        contract_group = df.groupby("inventory_contract").agg({
            "advertiser_cost_(adv_currency)": "sum",
            "tv_quality_index_raw": "sum",
            "tv_quality_index_measured_impressions": "sum"
        }).reset_index()
        contract_group["tvqi_score"] = contract_group["tv_quality_index_raw"] / contract_group["tv_quality_index_measured_impressions"]
        contract_group = contract_group.sort_values(by="advertiser_cost_(adv_currency)", ascending=False).head(10)

        fig2, ax3 = plt.subplots(figsize=(7, 3.5))
        ax4 = ax3.twinx()
        ax3.bar(contract_group["inventory_contract"], contract_group["advertiser_cost_(adv_currency)"], color="salmon")
        ax4.plot(contract_group["inventory_contract"], contract_group["tvqi_score"], color="darkblue", marker="o")
        ax3.set_ylabel("Advertiser Cost ($)", color="darkred")
        ax4.set_ylabel("TVQI Score", color="darkblue")
        ax3.set_title("By Inventory Contract", fontsize=12)
        ax3.tick_params(axis="x", rotation=45)
        ax3.yaxis.set_major_formatter(mtick.StrMethodFormatter("${x:,.0f}"))
        fig2.tight_layout()
        st.pyplot(fig2)

    # === Export Buttons ===
    st.subheader("📤 Export Visuals")
    col_pdf, col_ppt = st.columns(2)

    with col_pdf:
        if st.button("📄 Export Charts to PDF"):
            pdf_buf = io.BytesIO()
            fig1.savefig(pdf_buf, format="pdf")
            fig2.savefig(pdf_buf, format="pdf")
            st.download_button("Download PDF", pdf_buf.getvalue(), file_name="tvqi_charts.pdf", mime="application/pdf")

    with col_ppt:
        if st.button("📊 Export Charts to PowerPoint"):
            ppt = Presentation()
            for fig in [fig1, fig2]:
                slide = ppt.slides.add_slide(ppt.slide_layouts[5])
                image_stream = io.BytesIO()
                fig.savefig(image_stream, format="png")
                image_stream.seek(0)
                slide.shapes.add_picture(image_stream, Inches(1), Inches(1.5), Inches(8), Inches(4.5))
            ppt_bytes = io.BytesIO()
            ppt.save(ppt_bytes)
            st.download_button("Download PPT", ppt_bytes.getvalue(), file_name="tvqi_charts.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")

else:
    st.info("👈 Upload a CSV file to begin.")
